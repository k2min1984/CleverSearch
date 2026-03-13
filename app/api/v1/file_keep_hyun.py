import io
import re
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
from datetime import datetime

# FastAPI 관련 항목들을 한 줄로 통합
from fastapi import APIRouter, UploadFile, File, HTTPException, Query 
from fastapi.concurrency import run_in_threadpool

# 프로젝트 내부 모듈
from app.core.opensearch import get_client
from app.core.file import hwp, image
from app.common.utils import DocumentUtils
from app.common.embedding import embedder  #[추가] AI 벡터 변환기 가져오기

router = APIRouter()
client = get_client()
INDEX_NAME = "cleversearch-docs"
MAX_UPLOAD_BYTES = 20 * 1024 * 1024  # 20MB
EMBEDDING_DIMENSION = 768
_index_checked = False

# 허용된 확장자 목록 (doc, ppt는 라이브러리 미지원으로 제외)
ALLOWED_EXTENSIONS = {'xlsx', 'xls', 'hwp', 'hwpx', 'pdf', 'docx', 'pptx', 'jpg', 'jpeg', 'png'}


def _strip_control_chars(text: str) -> str:
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)


def _extract_text_content(content: bytes, ext: str) -> str:
    """확장자별 본문 텍스트 추출 (동기 함수, threadpool에서 실행)."""
    if ext in ['hwp', 'hwpx']:
        raw_text = hwp._extract_text(content, ext)
        return f"[[Page 1]]\n{_strip_control_chars(raw_text)}"

    if ext == 'pdf':
        with pdfplumber.open(io.BytesIO(content)) as pdf_file:
            pages = [
                f"[[Page {i+1}]]\n{p.extract_text()}"
                for i, p in enumerate(pdf_file.pages)
                if p.extract_text()
            ]
            return "\n\n".join(pages)

    if ext == 'docx':
        doc = Document(io.BytesIO(content))
        return "[[Page 1]]\n" + "\n".join([p.text for p in doc.paragraphs])

    if ext == 'pptx':
        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_txt = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            slides.append(f"[[Page {i+1}]]\n{slide_txt}")
        return "\n\n".join(slides)

    if ext in ['xlsx', 'xls']:
        excel_data = pd.read_excel(io.BytesIO(content), sheet_name=None)
        sheets = [f"[[Sheet: {name}]]\n{df.to_string()}" for name, df in excel_data.items()]
        return "\n\n".join(sheets)

    if ext in ['jpg', 'jpeg', 'png']:
        return image._extract_text(content)

    return ""

def ensure_index():
    """인덱스가 없을 경우 초기 설정(Settings/Mappings)과 함께 생성"""
    global _index_checked
    if _index_checked:
        return

    if not client.indices.exists(index=INDEX_NAME):
        # 🛡️ 여기가 우리가 다시 잡아야 할 '원점' 설계도입니다.
        index_body = {
            "settings": {
                "index": {
                    "knn": True,
                    "analysis": {
                        "analyzer": {
                            "korean_analyzer": { # 형태소 분석기 이름
                                "type": "custom",
                                "tokenizer": "nori_tokenizer",
                                "filter": ["lowercase", "nori_readingform", "my_stop_filter"]
                            }
                        },
                        "filter": {
                            "my_stop_filter": { # 불용어(쓰레기 단어) 처리
                                "type": "stop",
                                "stopwords": ["에", "대해서", "해주세요", "알려주세요", "의", "를", "은", "는"]
                            }
                        }
                    }
                }
            },
            "mappings": {
                "properties": {
                    # 🔥 [가장 중요] 기존 "korean" 대신 위에서 정의한 "korean_analyzer"를 써야 합니다.
                    "Title": {"type": "text", "analyzer": "korean_analyzer"},
                    "all_text": {"type": "text", "analyzer": "korean_analyzer"},
                    "doc_category": {"type": "keyword"},
                    "chosung_text": {"type": "text", "analyzer": "whitespace"},
                    "origin_file": {"type": "keyword"},
                    "file_ext": {"type": "keyword"},
                    "content_hash": {"type": "keyword"},
                    "indexed_at": {"type": "date"},
                    "text_vector": {
                        "type": "knn_vector",
                        "dimension": 768
                    }
                }
            }
        }
        
        try:
            client.indices.create(
                index=INDEX_NAME, 
                body=index_body
            )
            print(f"✅ 인덱스 [{INDEX_NAME}] 생성 성공")
        except Exception as e:
            raise RuntimeError(f"인덱스 생성 실패: {str(e)}")

    _index_checked = True

@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """
    파일 업로드 -> 텍스트 추출 -> 방역 -> 중복체크 -> 카테고리 분류 -> 저장 
    전체 프로세스를 담당하는 엔드포인트
    """
    # 1. 인덱스 존재 여부 확인 및 생성
    await run_in_threadpool(ensure_index)
    
    # 2. 파일명 방역 및 정보 추출
    filename = DocumentUtils.sanitize_text(file.filename or "")
    if not filename:
        raise HTTPException(status_code=400, detail="유효한 파일명이 필요합니다.")

    ext = filename.split('.')[-1].lower() if '.' in filename else ""
    
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"지원하지 않는 확장자입니다: {ext}")
        
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="빈 파일은 업로드할 수 없습니다.")
    if len(content) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail=f"파일 크기 제한({MAX_UPLOAD_BYTES // (1024 * 1024)}MB)을 초과했습니다.")

    try:
        # 3. 확장자별 텍스트 추출 (페이지 마커 삽입)
        text_content = await run_in_threadpool(_extract_text_content, content, ext)

        # 4. 공통 유틸리티를 이용한 본문 방역 (제어문자 제거)
        clean_body_text = DocumentUtils.sanitize_text(text_content)
        
        if not clean_body_text.strip():
            return {"status": "fail", "message": "추출된 텍스트가 없습니다."}

        # 5. 중복 확인용 지문(Digest) 생성
        content_digest = await run_in_threadpool(DocumentUtils.generate_content_digest, clean_body_text, filename)

        # 6. 내용 중복 체크 (DB 조회)
        is_dup, existing_file = await run_in_threadpool(
            DocumentUtils.check_duplicate_content,
            client,
            INDEX_NAME,
            content_digest,
        )
        if is_dup:
            return {
                "status": "skipped", 
                "message": f"내용 중복: 이미 '{existing_file}' 파일로 등록된 내용입니다."
            }

        # 7. 파일명 기반 카테고리 자동 분류 (PLAN, REPORT, RULE 등)
        category = DocumentUtils.map_category(filename)

        #최종 저장 전, 모든 필드의 텍스트에서 불법 제어문자(Code 31 등)를 싹 비웁니다.
        clean_body_text = _strip_control_chars(clean_body_text)
        # [추가] AI 두뇌로 본문 텍스트를 768차원 숫자로 변환
        # 주의: AI 모델은 한 번에 너무 긴 글을 읽으면 뻗으므로, 핵심 내용이 있는 앞부분 2000자만 잘라서 벡터로 만듭니다.
        vector_data = await run_in_threadpool(embedder.get_embedding, clean_body_text[:2000])
        if not isinstance(vector_data, list) or len(vector_data) != EMBEDDING_DIMENSION:
            raise HTTPException(status_code=500, detail="임베딩 벡터 차원이 올바르지 않습니다.")

        # 8. 최종 데이터 구조 생성
        doc_source = {
            "origin_file": str(filename),
            "all_text": str(clean_body_text),
            "doc_category": category,
            "chosung_text": DocumentUtils.convert_to_chosung(str(clean_body_text)),
            "content_hash": str(content_digest),
            "Title": str(filename),
            "file_ext": ext,
            "indexed_at": datetime.now().isoformat(),
            "text_vector": vector_data  #[추가] 768개의 숫자 리스트를 전용 방에 저장!
        }

        # 9. OpenSearch 저장
        await run_in_threadpool(client.index, index=INDEX_NAME, body=doc_source, refresh=True)

        return {
            "status": "success", 
            "message": f"[{filename}] 업로드 완료", 
            "category": category
        }

    except HTTPException:
        raise
    except Exception as e:
        # 에러 메시지도 방역 처리하여 전달
        clean_err = DocumentUtils.sanitize_text(str(e))
        raise HTTPException(status_code=500, detail=f"서버 내부 오류: {clean_err}")

# --- 관리용 API ---
@router.delete("/clear-all-data")
async def clear_all_data():
    client.delete_by_query(index=INDEX_NAME, body={"query": {"match_all": {}}}, refresh=True)
    return {"message": "전체 삭제 완료"}

@router.delete("/delete-index")
async def delete_index():
    if client.indices.exists(index=INDEX_NAME):
        client.indices.delete(index=INDEX_NAME)
    return {"message": "인덱스 완전 삭제 완료"}

@router.get("/search")
async def search_documents(keyword: str = Query(...)):
    """
    [하이브리드 검색] 키워드 매칭(전통 방식)과 AI 벡터 유사도(문맥 방식)를 동시에 검색합니다.
    """
    # 1. 검색어 방역 (특수문자 제거)
    safe_kw = re.sub(r'[^가-힣ㄱ-ㅎㅏ-ㅣa-zA-Z0-9\s]', '', keyword)
    if not safe_kw:
        return []

    # 2. 사용자의 검색어도 AI 두뇌를 거쳐 768차원 숫자(벡터)로 변환!
    query_vector = embedder.get_embedding(safe_kw)

    # 3. 궁극의 하이브리드 쿼리 조립 (AND 조건 타파 & AI 가중치 뻥튀기)
    hybrid_query = {
        "size": 10,
        "query": {
            "bool": {
                "should": [
                    {
                        "multi_match": {
                            "query": safe_kw,
                            "fields": ["Title^3", "all_text"],
                            "operator": "or"
                        }
                    },
                    {
                        "knn": {
                            "text_vector": {
                                "vector": query_vector,
                                "k": 10,
                                "boost": 1.0
                            }
                        }
                    }
                ],
                "minimum_should_match": 1
            }
        },
        # 🔥 [이게 정답입니다!] 0.5가 아니라 0.75로 세팅해야 유령 단어가 죽습니다!
        "min_score": 0.75, 
        
        "_source": {"excludes": ["text_vector"]},
        "highlight": {"fields": {"all_text": {}}}
    }

    # 4. OpenSearch에 하이브리드 쿼리 발사!
    res = client.search(index=INDEX_NAME, body=hybrid_query)

    # 🚀 [여기에 딱 2줄 추가!] 백엔드에서 실제로 몇 점을 주고 있는지 멱살 잡고 확인하기
    if res['hits']['hits']:
        print(f"👉 [디버그] 검색어: '{keyword}' / 1등 문서 점수: {res['hits']['hits'][0]['_score']}")

    # 검색된 문서 리스트 반환
    return res['hits']['hits']

@router.get("/read/{doc_id}")
async def read_document_full_text(doc_id: str):
    res = client.get(index=INDEX_NAME, id=doc_id)
    return res['_source']
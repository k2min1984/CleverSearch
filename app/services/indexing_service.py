"""
########################################################
# Description
# 색인 서비스 (IndexingService)
# 파일 파싱 → 텍스트 추출 → 청크 분할 → AI 임베딩 → OpenSearch 색인
# - PDF/HWP/HWPX/DOCX/PPTX/XLSX/이미지 파싱 통합
# - 텍스트 청크 분할 (chunk_size 단위)
# - 초성 텍스트 변환 (chosung_text 필드 생성)
# - 768차원 벡터 임베딩 (ko-sroberta-multitask)
# - OpenSearch 벨크 색인
#
# Modified History
# 강광묵 / 2026-01-20 / 최초생성
# 강광민 / 2026-02-13 / 파서 통합 및 임베딩 고도화
# 강광민 / 2026-03-23 / 헤더 주석 추가
########################################################
"""
import io
import json
import re
from datetime import datetime, timezone

import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation

from app.common.embedding import embedder
from app.common.utils import DocumentUtils
from app.core.file import hwp, image
from app.core.opensearch import get_client
from app.services.db_service import DBService


INDEX_NAME = "cleversearch-docs"
ALLOWED_EXTENSIONS = {"xlsx", "xls", "hwp", "hwpx", "pdf", "docx", "pptx", "jpg", "jpeg", "png", "txt"}


class IndexingService:
    @staticmethod
    def ensure_index() -> None:
        client = get_client()
        if client.indices.exists(index=INDEX_NAME):
            return

        index_body = {
            "settings": {
                "index": {
                    "knn": True,
                    "analysis": {
                        "analyzer": {
                            "korean_analyzer": {
                                "type": "custom",
                                "tokenizer": "nori_tokenizer",
                                "filter": ["lowercase", "nori_readingform", "my_stop_filter"],
                            }
                        },
                        "filter": {
                            "my_stop_filter": {
                                "type": "stop",
                                "stopwords": ["에", "대해서", "해주세요", "알려주세요", "의", "를", "은", "는"],
                            }
                        },
                    },
                }
            },
            "mappings": {
                "properties": {
                    "Title": {"type": "text", "analyzer": "korean_analyzer"},
                    "all_text": {"type": "text", "analyzer": "korean_analyzer"},
                    "doc_category": {"type": "keyword"},
                    "chosung_text": {"type": "text", "analyzer": "whitespace"},
                    "origin_file": {"type": "keyword"},
                    "file_ext": {"type": "keyword"},
                    "content_hash": {"type": "keyword"},
                    "indexed_at": {"type": "date"},
                    "text_vector": {"type": "knn_vector", "dimension": 768},
                }
            },
        }
        client.indices.create(index=INDEX_NAME, body=index_body)

    @staticmethod
    def _extract_text(filename: str, content: bytes) -> tuple[str, str]:
        ext = filename.split(".")[-1].lower() if "." in filename else ""
        if ext not in ALLOWED_EXTENSIONS:
            return "", ext

        text_content = ""
        if ext in ["hwp", "hwpx"]:
            raw_text = hwp._extract_text(content, ext)
            raw_text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", raw_text)
            text_content = f"[[Page 1]]\n{raw_text}"
        elif ext == "pdf":
            with pdfplumber.open(io.BytesIO(content)) as pdf_file:
                pages = [f"[[Page {i + 1}]]\n{p.extract_text()}" for i, p in enumerate(pdf_file.pages) if p.extract_text()]
                text_content = "\n\n".join(pages)

            if not text_content.strip():
                try:
                    import fitz

                    pdf_doc = fitz.open(stream=content, filetype="pdf")
                    ocr_pages = []
                    for page_num, page in enumerate(pdf_doc):
                        pix = page.get_pixmap(dpi=150)
                        img_bytes = pix.tobytes("png")
                        page_text = image._extract_text(img_bytes)
                        if page_text and page_text.strip():
                            ocr_pages.append(f"[[Page {page_num + 1}]]\n{page_text}")
                    text_content = "\n\n".join(ocr_pages)
                except Exception:
                    text_content = ""
        elif ext == "docx":
            doc = Document(io.BytesIO(content))
            text_content = f"[[Page 1]]\n" + "\n".join([p.text for p in doc.paragraphs])
        elif ext == "pptx":
            prs = Presentation(io.BytesIO(content))
            slides = []
            for i, slide in enumerate(prs.slides):
                slide_txt = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                slides.append(f"[[Page {i + 1}]]\n{slide_txt}")
            text_content = "\n\n".join(slides)
        elif ext in ["xlsx", "xls"]:
            excel_data = pd.read_excel(io.BytesIO(content), sheet_name=None)
            sheets = [f"[[Sheet: {name}]]\n{df.to_string()}" for name, df in excel_data.items()]
            text_content = "\n\n".join(sheets)
        elif ext in ["jpg", "jpeg", "png"]:
            text_content = image._extract_text(content)
        elif ext == "txt":
            text_content = content.decode("utf-8", errors="ignore")

        return text_content, ext

    @staticmethod
    def index_bytes(filename: str, content: bytes, source_label: str = "upload") -> dict:
        IndexingService.ensure_index()
        client = get_client()

        safe_filename = DocumentUtils.sanitize_text(filename)
        text_content, ext = IndexingService._extract_text(safe_filename, content)

        if ext not in ALLOWED_EXTENSIONS:
            return {"status": "fail", "message": f"지원하지 않는 확장자: {ext}", "file": safe_filename}

        clean_body_text = DocumentUtils.sanitize_text(text_content)
        clean_body_text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", clean_body_text)
        if not clean_body_text.strip():
            return {"status": "fail", "message": "추출된 텍스트 없음", "file": safe_filename}

        content_digest = DocumentUtils.generate_content_digest(clean_body_text, safe_filename)
        is_dup, existing_file = DocumentUtils.check_duplicate_content(client, INDEX_NAME, content_digest)
        if is_dup:
            return {
                "status": "skipped",
                "message": f"내용 중복: {existing_file}",
                "file": safe_filename,
                "content_hash": content_digest,
            }

        category = DocumentUtils.map_category(safe_filename)
        vector_data = embedder.get_embedding(clean_body_text[:2000])

        doc_source = {
            "origin_file": safe_filename,
            "all_text": clean_body_text,
            "doc_category": category,
            "chosung_text": DocumentUtils.convert_to_chosung(clean_body_text),
            "content_hash": str(content_digest),
            "Title": safe_filename,
            "file_ext": ext,
            "indexed_at": datetime.now(timezone.utc).isoformat(),
            "text_vector": vector_data,
            "source_label": source_label,
        }

        safe_doc = json.dumps(doc_source, ensure_ascii=False)
        index_res = client.index(index=INDEX_NAME, body=safe_doc, refresh=True)
        DBService.save_indexed_document(
            os_doc_id=index_res.get("_id", ""),
            origin_file=safe_filename,
            file_ext=ext,
            doc_category=category,
            content_hash=str(content_digest),
            title=safe_filename,
            all_text=clean_body_text,
        )

        return {
            "status": "success",
            "message": "색인 완료",
            "file": safe_filename,
            "content_hash": content_digest,
            "category": category,
            "doc_id": index_res.get("_id"),
        }

"""
########################################################
# Description
# Office 파일 파서 (Word / PowerPoint)
# DOCX, PPTX 문서에서 본문 텍스트 추출
# - Word: python-docx → 문단 + 테이블 + 헤더/푸터 추출
# - PowerPoint: python-pptx → 슬라이드 텍스트 + 테이블 + 노트 추출
#
# Modified History
# 강광묵 / 2026-01-20 / 최초생성
# 강광민 / 2026-04-02 / 고도화 (테이블, 노트, 헤더/푸터 추출 추가)
########################################################
"""
import io
from docx import Document
from pptx import Presentation


def parse_word(content: bytes) -> str:
    try:
        doc = Document(io.BytesIO(content))
        parts: list[str] = []

        # 본문 문단
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)

        # 테이블
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        # 헤더/푸터
        for section in doc.sections:
            for header_part in (section.header, section.first_page_header, section.even_page_header):
                if header_part and header_part.paragraphs:
                    for p in header_part.paragraphs:
                        if p.text.strip():
                            parts.append(p.text)
            for footer_part in (section.footer, section.first_page_footer, section.even_page_footer):
                if footer_part and footer_part.paragraphs:
                    for p in footer_part.paragraphs:
                        if p.text.strip():
                            parts.append(p.text)

        return "[[Page 1]]\n" + "\n".join(parts)
    except Exception:
        return ""


def parse_pptx(content: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(content))
        slides: list[str] = []

        for i, slide in enumerate(prs.slides):
            parts: list[str] = []

            # 셰이프 텍스트
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)

                # 테이블
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))

            # 슬라이드 노트
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(notes)

            if parts:
                slides.append(f"[[Page {i + 1}]]\n" + "\n".join(parts))

        return "\n\n".join(slides)
    except Exception:
        return ""


def extract_text(content: bytes, ext: str) -> str:
    """DOCX/PPTX 통합 진입점. 확장자로 분기."""
    if ext == "docx":
        return parse_word(content)
    elif ext == "pptx":
        return parse_pptx(content)
    return ""